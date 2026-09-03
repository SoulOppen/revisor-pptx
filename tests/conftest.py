"""Shared pytest fixtures for revisor-pptx.

The fixture factory builds real .pptx files programmatically (never static
files), so tests exercise the actual python-pptx I/O path.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def make_pptx(tmp_path):
    """Factory: build a .pptx with controlled text and save it under tmp_path.

    Constructs slides with a text shape, a table, and speaker notes. Returns a
    callable ``make_pptx(name="test.pptx", with_table=True, with_notes=True,
    texts=None)`` that produces the file path.
    """

    def _make(
        name: str = "test.pptx",
        with_table: bool = True,
        with_notes: bool = True,
        texts: dict[int, str] | None = None,
    ) -> Path:
        prs = Presentation()
        default_texts = {1: "hola munod como estas", 2: "tengo una problema"}
        texts = texts or default_texts

        # Slide 1: a text shape with a deliberate typo.
        slide1 = prs.slides.add_slide(prs.slide_layouts[1])
        box1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf1 = box1.text_frame
        tf1.text = texts.get(1, "")
        if with_table:
            rows, cols = 1, 1
            left, top, width, height = Inches(1), Inches(2), Inches(3), Inches(1)
            table_shape = slide1.shapes.add_table(rows, cols, left, top, width, height)
            table_shape.table.cell(0, 0).text = "comprar leche"
        if with_notes:
            slide1.notes_slide.notes_text_frame.text = "nota con errata final"

        # Slide 2: notes only.
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        if with_notes:
            slide2.notes_slide.notes_text_frame.text = "segunda nota sin error"
        box2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box2.text_frame.text = texts.get(2, "")

        # Slide 3: empty (no text at all).
        prs.slides.add_slide(prs.slide_layouts[6])

        path = tmp_path / name
        prs.save(path)
        return path

    return _make
