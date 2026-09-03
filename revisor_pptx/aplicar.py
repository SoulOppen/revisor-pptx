"""High-confidence correction filtering and application to the copy.

``filter_corrections`` is PURE (unit tested). ``apply_corrections`` is the IO
boundary that mutates the copied .pptx file, preserving the formatting of
unaffected runs.
"""

from __future__ import annotations

from collections.abc import Iterable

from .extract_text import SlideText
from .revisar import Correction

# Issue types we actively fix.
_INCLUDE_ISSUES = {"misspelling", "grammar", "typo", "typos", "inconsistency"}
# Issue types that are never applied automatically.
_EXCLUDE_ISSUES = {"style", "whitespace", "casing"}


def filter_corrections(
    corrections: Iterable[Correction], replacements_limit: int = 3
) -> list[Correction]:
    """Pure: keep high-confidence corrections to auto-apply.

    A correction is applied whenever it is an actionable issue type with at
    least one replacement. The first (best) replacement is written into the
    document; any other options are only recorded in the report.
    """
    kept: list[Correction] = []
    for c in corrections:
        issue = c.rule_issue
        if issue in _EXCLUDE_ISSUES:
            continue
        if issue not in _INCLUDE_ISSUES:
            continue

        replacements = [r for r in c.replacements if isinstance(r, str) and r]
        if not replacements:
            continue

        kept.append(c)
    return kept


def _best_replacement(c: Correction) -> str:
    """The replacement to write into the document: the first (best) option.

    LanguageTool orders replacements by likelihood, so the first entry is the
    preferred fix. The remaining options are still recorded in the report.
    """
    if not c.replacements:
        return ""
    return c.replacements[0]


def apply_corrections(
    path, slides: list[SlideText], corrections: list[Correction]
) -> list[Correction]:
    """IO boundary: apply corrections to the copied .pptx at ``path``.

    Re-opens the file, locates the matching text by offsets, and replaces
    preserving formatting where possible. Returns the list of corrections that
    were actually applied.
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    applied_corrections: list[Correction] = []

    for slide_text, slide_obj in zip(slides, prs.slides):
        slide_corrs = [c for c in corrections if c.slide_idx == slide_text.slide_idx]

        # Group by shape: iterate shapes in order matching their shape_idx.
        for st in slide_text.shapes_text:
            shape_corrs = [
                c
                for c in slide_corrs
                if c.shape_idx == st.shape_idx and c.shape_idx != -1
            ]
            shape_obj = _find_shape(slide_obj, st.shape_idx)
            if shape_obj is None or not shape_corrs:
                continue

            if getattr(shape_obj, "has_table", False):
                # Per-cell: rebuild cumulative offsets and apply within each cell text frame.
                applied_corrections += _apply_table_cells(
                    shape_obj.table, shape_corrs, st.segments
                )
            else:
                applied_corrections += _apply_text_frame(
                    shape_obj.text_frame, shape_corrs
                )

        # Notes corrections are a single source; routed via shape_idx == -1.
        notes_corrs = [c for c in slide_corrs if c.shape_idx == -1]
        if notes_corrs and slide_obj.has_notes_slide:
            frame = slide_obj.notes_slide.notes_text_frame
            applied_corrections += _apply_text_frame(frame, notes_corrs)

    prs.save(str(path))
    return applied_corrections


def _find_shape(slide_obj, shape_idx):
    for shape in slide_obj.shapes:
        if shape.shape_id == shape_idx:
            return shape
    return None


def _apply_text_frame(frame, corrections) -> list[Correction]:
    """Apply corrections within a single TextFrame by replacing run text.

    Offsets in each correction are relative to the concatenated text of the
    frame's runs; this rebuilds that concatenation cumulatively to locate the
    matching run and replaces preserving the run's formatting.
    """
    paragraphs = list(frame.paragraphs)
    applied: list[Correction] = []
    for c in sorted(corrections, key=lambda x: x.offset):
        cumulative = 0
        for p in paragraphs:
            for r in p.runs:
                seg_start = cumulative
                seg_end = cumulative + len(r.text)
                if seg_start <= c.offset < seg_end:
                    inner = c.offset - seg_start
                    if r.text[inner : inner + c.length] == c.original:
                        r.text = (
                            r.text[:inner]
                            + _best_replacement(c)
                            + r.text[inner + c.length :]
                        )
                        applied.append(c)
                    break
                cumulative = seg_end
    return applied


def _apply_table_cells(table, corrections, segments) -> list[Correction]:
    applied: list[Correction] = []
    # Recompute a flat list of cells and their cumulative text offsets.
    cells = [(row, cell) for row in table.rows for cell in row.cells]
    cumulative = 0
    cell_offsets: list[tuple] = []
    for row, cell in cells:
        text = "".join(run.text for p in cell.text_frame.paragraphs for run in p.runs)
        cell_offsets.append((text, cumulative, row, cell))
        cumulative += len(text)

    for c in sorted(corrections, key=lambda x: x.offset):
        for text, start, row, cell in cell_offsets:
            if start <= c.offset < start + len(text):
                inner = c.offset - start
                if text[inner : inner + c.length] == c.original:
                    # Apply within the cell's own text frame.
                    _apply_whole_cell(
                        cell.text_frame, inner, c.length, _best_replacement(c)
                    )
                    applied.append(c)
                break
    return applied


def _apply_whole_cell(text_frame, inner, length, replacement) -> None:
    """Replace text inside a cell text frame at a run, preserving run format."""
    cumulative = 0
    for p in text_frame.paragraphs:
        for r in p.runs:
            if cumulative <= inner < cumulative + len(r.text):
                local_inner = inner - cumulative
                r.text = (
                    r.text[:local_inner] + replacement + r.text[local_inner + length :]
                )
                return
            cumulative += len(r.text)
