"""Pure text extraction from presentation slides.

The pure functions here take pptx objects (python-pptx) and return plain data
structures. No disk access happens inside the pure layer.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

SegmentSource = Literal["shape", "table", "notes"]


@dataclass
class TextSegment:
    """A contiguous block of text with its cumulative offset within a shape.

    ``offset`` is the character position where this segment starts relative to
    the start of its containing shape's concatenated text.
    """

    text: str
    offset: int
    source: SegmentSource


@dataclass
class ShapeText:
    """All text belonging to a single shape (or table cell run group)."""

    shape_idx: int
    shape_name: str
    segments: list[TextSegment] = field(default_factory=list)


@dataclass
class SlideText:
    """All extractable text of a single slide."""

    slide_idx: int
    shapes_text: list[ShapeText] = field(default_factory=list)
    notes: str = ""


def _segments_from_text_frame(text_frame) -> list[TextSegment]:
    """Build text segments from a TextFrame, preserving cumulative offsets."""
    segments: list[TextSegment] = []
    total = 0
    paragraphs = getattr(text_frame, "paragraphs", [])
    for paragraph in paragraphs:
        text = "".join(run.text for run in paragraph.runs)
        if text:
            segments.append(TextSegment(text=text, offset=total, source="shape"))
            total += len(text)
    return segments


def _segments_from_table(table) -> list[TextSegment]:
    """Build one ascending-offset sequence per table cell text frame."""
    segments: list[TextSegment] = []
    cumulative = 0
    for row in table.rows:
        for cell in row.cells:
            text_frame = cell.text_frame
            text = "".join(
                run.text
                for paragraph in text_frame.paragraphs
                for run in paragraph.runs
            )
            if text:
                segments.append(
                    TextSegment(text=text, offset=cumulative, source="table")
                )
                cumulative += len(text)
    return segments


def extract_slide_text(slide) -> SlideText:
    """Pure: extract all visible text from a single slide object.

    Covers shapes (text frames), tables (cell by cell), and speaker notes.
    Empty text blocks are silently skipped.
    """
    slide_text = SlideText(slide_idx=slide.slide_id)

    for shape in slide.shapes:
        shape_idx = shape.shape_id
        shape_name = shape.name
        segments: list[TextSegment] = []

        if getattr(shape, "has_table", False):
            segments += _segments_from_table(shape.table)
            slide_text.shapes_text.append(
                ShapeText(shape_idx=shape_idx, shape_name=shape_name, segments=segments)
            )
            continue

        if hasattr(shape, "text_frame") and shape.has_text_frame:
            segments += _segments_from_text_frame(shape.text_frame)
            slide_text.shapes_text.append(
                ShapeText(shape_idx=shape_idx, shape_name=shape_name, segments=segments)
            )

    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        notes = "".join(
            run.text for paragraph in notes_frame.paragraphs for run in paragraph.runs
        )
        slide_text.notes = notes

    return slide_text


def extract_pptx(path: Path) -> list[SlideText]:
    """IO boundary: open a .pptx from disk and extract text per slide."""
    from pptx import Presentation

    prs = Presentation(str(path))
    return [extract_slide_text(slide) for slide in prs.slides]
